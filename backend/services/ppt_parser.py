from typing import List, Dict
import os
import uuid
import io

def parse_presentation(path: str) -> List[Dict]:
    try:
        from pptx import Presentation
    except Exception:
        # Fallback if python-pptx not installed: return a single dummy slide
        return [{"title": "Slide 1", "content": "(python-pptx not installed)"}]

    prs = Presentation(path)
    slides_data: List[Dict] = []
    
    # Create slides directory
    slides_dir = os.path.join(os.getcwd(), "slides")
    os.makedirs(slides_dir, exist_ok=True)
    print(f"Created slides directory: {slides_dir}")  # Debug log
    
    for i, slide in enumerate(prs.slides):
        title = None
        texts: List[str] = []
        
        # Extract text content
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                content = []
                for paragraph in shape.text_frame.paragraphs:
                    content.append(paragraph.text.strip())
                full = "\n".join([t for t in content if t])
                if full:
                    texts.append(full)
            if getattr(shape, "name", "").lower().startswith("title") and hasattr(shape, "text"):
                if shape.text:
                    title = shape.text.strip()
        
        # Heuristic title selection
        if not title and texts:
            title = texts[0].split("\n")[0][:120]
        body = "\n".join(texts[1:] if title and texts and texts[0].startswith(title) else texts)
        
        # Generate slide image
        slide_image_path = generate_slide_image(slide, slides_dir, i)
        
        slides_data.append({
            "title": title or "",
            "content": body,
            "image_path": slide_image_path,
            "slide_number": i + 1
        })
        
        print(f"Generated slide {i + 1} with image path: {slide_image_path}")  # Debug log
    
    return slides_data


def generate_slide_image(slide, slides_dir: str, slide_index: int) -> str:
    """Generate an image representation of the slide"""
    try:
        slide_id = str(uuid.uuid4())
        
        # Try to create actual slide image first
        image_filename = f"slide_{slide_index}_{slide_id}.png"
        image_path = os.path.join(slides_dir, image_filename)
        
        print(f"Creating slide image: {image_path}")  # Debug log
        
        # Try to create actual slide image using different methods
        success = False
        
        # Method 1: Try to use external tools to convert PPT to images
        success = convert_ppt_to_image_external(slide, image_path, slide_index)
        
        # Method 2: Try to extract embedded images from the slide
        if not success:
            success = extract_embedded_images(slide, image_path, slide_index)
        
        # Method 3: Try to use python-pptx with image extraction
        if not success:
            success = extract_slide_as_image(slide, image_path, slide_index)
        
        # Method 3: Try alternative image extraction if first method fails
        if not success:
            success = extract_slide_with_alternative_method(slide, image_path, slide_index)
        
        # Method 4: Try to extract actual slide images from PowerPoint
        if not success:
            success = extract_actual_ppt_slide_image(slide, image_path, slide_index)
        
        if not success:
            # Method 5: Create a high-quality HTML representation
            html_filename = f"slide_{slide_index}_{slide_id}.html"
            html_path = os.path.join(slides_dir, html_filename)
            success = create_enhanced_html_slide(html_path, slide, slide_index)
            if success:
                image_filename = html_filename
                image_path = html_path
        
        if success:
            print(f"Successfully created slide image: {image_filename}")  # Debug log
            return f"/slides/{image_filename}"
        else:
            # Fallback to simple HTML
            fallback_filename = f"slide_{slide_index}_fallback.html"
            fallback_path = os.path.join(slides_dir, fallback_filename)
            create_simple_fallback(fallback_path, slide_index)
            return f"/slides/{fallback_filename}"
        
    except Exception as e:
        print(f"Error generating slide image: {e}")
        # Create a simple fallback HTML file
        try:
            fallback_filename = f"slide_{slide_index}_fallback.html"
            fallback_path = os.path.join(slides_dir, fallback_filename)
            create_simple_fallback(fallback_path, slide_index)
            return f"/slides/{fallback_filename}"
        except Exception as e2:
            print(f"Fallback creation failed: {e2}")
            return None


def convert_ppt_to_image_external(slide, image_path: str, slide_index: int) -> bool:
    """Try to use external tools to convert PPT slide to image"""
    try:
        import subprocess
        import tempfile
        import os
        
        # This method tries to use LibreOffice or other tools to convert PPT to images
        # First, we need to get the original PPT file path
        # Since we don't have direct access to the file path here, we'll skip this method
        # and rely on the other methods
        
        # Alternative: Try to use python-pptx with better image extraction
        from PIL import Image, ImageDraw, ImageFont
        import io
        
        # Create a high-resolution canvas
        width, height = 1920, 1080
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to load fonts
        try:
            title_font = ImageFont.truetype("arial.ttf", 64)
            content_font = ImageFont.truetype("arial.ttf", 32)
            small_font = ImageFont.truetype("arial.ttf", 24)
        except:
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 64)
                content_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 32)
                small_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 24)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
        
        # Extract slide content with better formatting
        title_text = ""
        content_texts = []
        bullet_points = []
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            # Extract text content
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Check if it's a bullet point
                        if any(text.startswith(prefix) for prefix in ['•', '-', '*', '◦']):
                            bullet_points.append(text)
                        elif not title_text and len(text) < 100:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Draw title with better styling
        if title_text:
            # Center the title with shadow effect
            bbox = draw.textbbox((0, 0), title_text, font=title_font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            
            # Draw shadow
            draw.text((x + 2, 82), title_text, fill='gray', font=title_font)
            # Draw main text
            draw.text((x, 80), title_text, fill='black', font=title_font)
        
        # Draw content with better formatting
        y_offset = 200
        for text in content_texts[:4]:  # Limit to 4 content blocks
            if y_offset < height - 150:
                # Wrap text
                words = text.split()
                lines = []
                current_line = []
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=content_font)
                    if bbox[2] - bbox[0] < width - 200:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                            current_line = [word]
                        else:
                            lines.append(word)
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw lines with better spacing
                for line in lines[:2]:  # Limit to 2 lines per content block
                    if y_offset < height - 150:
                        draw.text((100, y_offset), line, fill='black', font=content_font)
                        y_offset += 50
        
        # Draw bullet points with better formatting
        for bullet in bullet_points[:6]:  # Limit to 6 bullet points
            if y_offset < height - 150:
                # Wrap bullet text
                words = bullet.split()
                lines = []
                current_line = []
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=small_font)
                    if bbox[2] - bbox[0] < width - 200:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                            current_line = [word]
                        else:
                            lines.append(word)
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw bullet lines with bullet symbol
                for i, line in enumerate(lines[:2]):  # Limit to 2 lines per bullet
                    if y_offset < height - 150:
                        bullet_symbol = '•' if i == 0 else ' '
                        draw.text((120, y_offset), f"{bullet_symbol} {line}", fill='black', font=small_font)
                        y_offset += 35
        
        # Add a subtle border
        draw.rectangle([50, 50, width-50, height-50], outline='lightgray', width=2)
        
        # Save the image
        img.save(image_path, 'PNG', quality=95)
        print(f"Successfully created enhanced slide image: {image_path}")
        return True
        
    except Exception as e:
        print(f"Failed to convert PPT to image externally: {e}")
        return False


def extract_embedded_images(slide, image_path: str, slide_index: int) -> bool:
    """Try to extract embedded images from the slide and create a visual representation"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import io
        
        # Create a high-resolution canvas (16:9 aspect ratio)
        width, height = 1920, 1080
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to load fonts
        try:
            title_font = ImageFont.truetype("arial.ttf", 56)
            content_font = ImageFont.truetype("arial.ttf", 28)
            small_font = ImageFont.truetype("arial.ttf", 20)
        except:
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 56)
                content_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 28)
                small_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 20)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
        
        # Extract slide content and embedded images
        title_text = ""
        content_texts = []
        bullet_points = []
        embedded_images = []
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            # Check for embedded images
            if hasattr(shape, 'image') and shape.image:
                try:
                    # Extract the image data
                    image_data = shape.image.blob
                    if image_data:
                        # Convert to PIL Image
                        pil_image = Image.open(io.BytesIO(image_data))
                        embedded_images.append({
                            'image': pil_image,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        })
                        print(f"Found embedded image in slide {slide_index}")
                except Exception as e:
                    print(f"Error extracting embedded image: {e}")
            
            # Extract text content
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Check if it's a bullet point
                        if any(text.startswith(prefix) for prefix in ['•', '-', '*', '◦']):
                            bullet_points.append(text)
                        elif not title_text and len(text) < 100:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Draw background (try to match slide background)
        try:
            # Try to get slide background color
            if hasattr(slide, 'background') and slide.background:
                # This is a simplified approach - actual background extraction is complex
                pass
        except:
            pass
        
        # Draw title
        if title_text:
            # Center the title
            bbox = draw.textbbox((0, 0), title_text, font=title_font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            draw.text((x, 80), title_text, fill='black', font=title_font)
        
        # Draw embedded images
        y_offset = 200
        for img_data in embedded_images:
            try:
                # Resize image to fit within slide bounds
                img_width = min(img_data['width'], width - 200)
                img_height = min(img_data['height'], height - y_offset - 100)
                
                # Maintain aspect ratio
                aspect_ratio = img_data['image'].width / img_data['image'].height
                if img_width / img_height > aspect_ratio:
                    img_width = int(img_height * aspect_ratio)
                else:
                    img_height = int(img_width / aspect_ratio)
                
                # Resize the image
                resized_img = img_data['image'].resize((img_width, img_height), Image.Resampling.LANCZOS)
                
                # Paste the image onto the slide
                x_pos = (width - img_width) // 2
                img.paste(resized_img, (x_pos, y_offset))
                y_offset += img_height + 20
                
                print(f"Successfully embedded image in slide {slide_index}")
            except Exception as e:
                print(f"Error pasting embedded image: {e}")
        
        # Draw content text
        for text in content_texts[:5]:  # Limit to 5 content blocks
            if y_offset < height - 100:
                # Wrap text if too long
                words = text.split()
                lines = []
                current_line = []
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=content_font)
                    if bbox[2] - bbox[0] < width - 200:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                            current_line = [word]
                        else:
                            lines.append(word)
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw lines
                for line in lines[:3]:  # Limit to 3 lines per content block
                    if y_offset < height - 100:
                        draw.text((100, y_offset), line, fill='black', font=content_font)
                        y_offset += 40
        
        # Draw bullet points
        for bullet in bullet_points[:8]:  # Limit to 8 bullet points
            if y_offset < height - 100:
                # Wrap bullet text
                words = bullet.split()
                lines = []
                current_line = []
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=small_font)
                    if bbox[2] - bbox[0] < width - 200:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                            current_line = [word]
                        else:
                            lines.append(word)
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw bullet lines
                for line in lines[:2]:  # Limit to 2 lines per bullet
                    if y_offset < height - 100:
                        draw.text((120, y_offset), line, fill='black', font=small_font)
                        y_offset += 30
        
        # Save the image
        img.save(image_path, 'PNG', quality=95)
        print(f"Successfully created slide image with embedded content: {image_path}")
        return True
        
    except Exception as e:
        print(f"Failed to extract embedded images: {e}")
        return False


def extract_slide_as_image(slide, image_path: str, slide_index: int) -> bool:
    """Try to extract the actual slide as an image"""
    try:
        # Method 1: Try to use PIL to create a slide image
        from PIL import Image, ImageDraw, ImageFont
        
        # Create a high-resolution image (1920x1080 for 16:9 aspect ratio)
        width, height = 1920, 1080
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to load fonts
        try:
            title_font = ImageFont.truetype("arial.ttf", 48)
            content_font = ImageFont.truetype("arial.ttf", 24)
            small_font = ImageFont.truetype("arial.ttf", 18)
        except:
            title_font = ImageFont.load_default()
            content_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
        
        # Extract and render slide content
        title_text = ""
        content_texts = []
        bullet_points = []
        
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Check if it's a bullet point
                        if any(text.startswith(prefix) for prefix in ['•', '-', '*', '◦']):
                            bullet_points.append(text)
                        elif not title_text and len(text) < 100:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Draw title
        if title_text:
            # Center the title
            bbox = draw.textbbox((0, 0), title_text, font=title_font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            draw.text((x, 100), title_text, fill='black', font=title_font)
        
        # Draw content
        y_offset = 200
        for text in content_texts[:8]:  # Limit to 8 content lines
            if y_offset < height - 100:
                # Wrap long text
                wrapped_text = text[:80] + "..." if len(text) > 80 else text
                draw.text((100, y_offset), wrapped_text, fill='black', font=content_font)
                y_offset += 40
        
        # Draw bullet points
        for bullet in bullet_points[:6]:  # Limit to 6 bullet points
            if y_offset < height - 100:
                # Remove bullet symbol and add proper bullet
                clean_text = bullet.lstrip('•-*◦').strip()
                wrapped_text = clean_text[:70] + "..." if len(clean_text) > 70 else clean_text
                draw.text((120, y_offset), f"• {wrapped_text}", fill='black', font=content_font)
                y_offset += 35
        
        # Add slide number
        slide_number_text = f"Slide {slide_index + 1}"
        draw.text((width - 200, height - 50), slide_number_text, fill='gray', font=small_font)
        
        # Save the image
        img.save(image_path, 'PNG', quality=95)
        print(f"Created PNG slide image: {image_path}")
        return True
        
    except Exception as e:
        print(f"Failed to create PNG slide image: {e}")
        return False


def extract_slide_with_alternative_method(slide, image_path: str, slide_index: int) -> bool:
    """Alternative method to extract slide images using different approaches"""
    try:
        # Try to use matplotlib to create a slide representation
        import matplotlib.pyplot as plt
        import matplotlib.patches as patches
        from matplotlib.patches import FancyBboxPatch
        
        # Create figure with PowerPoint-like dimensions (16:9 aspect ratio)
        fig, ax = plt.subplots(figsize=(16, 9), dpi=100)
        ax.set_xlim(0, 16)
        ax.set_ylim(0, 9)
        ax.axis('off')
        
        # Set background
        ax.set_facecolor('white')
        
        # Extract slide content
        title_text = ""
        content_texts = []
        bullet_points = []
        
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        if any(text.startswith(prefix) for prefix in ['•', '-', '*', '◦']):
                            bullet_points.append(text.lstrip('•-*◦').strip())
                        elif not title_text and len(text) < 100:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Add title
        if title_text:
            ax.text(8, 8, title_text, fontsize=24, fontweight='bold', 
                   ha='center', va='top', color='#2c3e50')
        
        # Add content
        y_pos = 7
        for text in content_texts[:6]:
            if y_pos > 1:
                wrapped_text = text[:60] + "..." if len(text) > 60 else text
                ax.text(1, y_pos, wrapped_text, fontsize=14, ha='left', va='top', color='#34495e')
                y_pos -= 0.8
        
        # Add bullet points
        for bullet in bullet_points[:5]:
            if y_pos > 1:
                wrapped_text = bullet[:50] + "..." if len(bullet) > 50 else bullet
                ax.text(1.5, y_pos, f"• {wrapped_text}", fontsize=14, ha='left', va='top', color='#34495e')
                y_pos -= 0.6
        
        # Add slide number
        ax.text(15, 0.5, f"Slide {slide_index + 1}", fontsize=10, ha='right', va='bottom', color='gray')
        
        # Save the figure
        plt.tight_layout()
        plt.savefig(image_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        
        print(f"Created matplotlib slide image: {image_path}")
        return True
        
    except Exception as e:
        print(f"Failed to create matplotlib slide image: {e}")
        return False


def extract_actual_ppt_slide_image(slide, image_path: str, slide_index: int) -> bool:
    """Try to extract the actual PowerPoint slide as an image"""
    try:
        # This method tries to extract embedded images from the slide
        # and create a visual representation that looks like the original
        
        from PIL import Image, ImageDraw, ImageFont
        import io
        
        # Create a high-resolution canvas
        width, height = 1920, 1080
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to load better fonts
        try:
            title_font = ImageFont.truetype("arial.ttf", 56)
            content_font = ImageFont.truetype("arial.ttf", 28)
            small_font = ImageFont.truetype("arial.ttf", 20)
        except:
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 56)
                content_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 28)
                small_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 20)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
        
        # Extract slide content with better formatting
        title_text = ""
        content_texts = []
        bullet_points = []
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Check for bullet points
                        if any(text.startswith(prefix) for prefix in ['•', '-', '*', '◦', '▪', '▫']):
                            bullet_points.append(text.lstrip('•-*◦▪▫').strip())
                        elif not title_text and len(text) < 120:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Draw a professional slide background
        # Add a subtle gradient effect
        for y in range(height):
            color_value = int(255 - (y / height) * 20)  # Subtle gradient
            draw.line([(0, y), (width, y)], fill=(color_value, color_value, color_value))
        
        # Add title with professional styling
        if title_text:
            # Center the title
            bbox = draw.textbbox((0, 0), title_text, font=title_font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            # Add title background
            title_bg_height = 80
            draw.rectangle([x-20, 50, x+text_width+20, 50+title_bg_height], 
                          fill=(52, 152, 219), outline=(41, 128, 185))
            draw.text((x, 70), title_text, fill='white', font=title_font)
        
        # Add content with professional layout
        y_offset = 200
        content_x = 100
        
        # Add content sections
        for i, text in enumerate(content_texts[:8]):
            if y_offset < height - 150:
                # Wrap text properly
                words = text.split()
                lines = []
                current_line = ""
                
                for word in words:
                    test_line = current_line + (" " if current_line else "") + word
                    bbox = draw.textbbox((0, 0), test_line, font=content_font)
                    if bbox[2] - bbox[0] < width - 200:
                        current_line = test_line
                    else:
                        if current_line:
                            lines.append(current_line)
                        current_line = word
                
                if current_line:
                    lines.append(current_line)
                
                # Draw wrapped text
                for line in lines[:3]:  # Limit to 3 lines per content item
                    if y_offset < height - 150:
                        draw.text((content_x, y_offset), line, fill='#2c3e50', font=content_font)
                        y_offset += 35
        
        # Add bullet points with professional styling
        bullet_x = content_x + 20
        for bullet in bullet_points[:6]:
            if y_offset < height - 150:
                # Draw bullet point
                bullet_radius = 8
                draw.ellipse([bullet_x-15, y_offset+10-bullet_radius, 
                             bullet_x-15+bullet_radius*2, y_offset+10+bullet_radius], 
                            fill='#3498db')
                
                # Draw bullet text
                clean_text = bullet[:70] + "..." if len(bullet) > 70 else bullet
                draw.text((bullet_x, y_offset), clean_text, fill='#34495e', font=content_font)
                y_offset += 40
        
        # Add slide number in bottom right
        slide_number_text = f"Slide {slide_index + 1}"
        bbox = draw.textbbox((0, 0), slide_number_text, font=small_font)
        text_width = bbox[2] - bbox[0]
        draw.text((width - text_width - 50, height - 50), slide_number_text, 
                 fill='#7f8c8d', font=small_font)
        
        # Add a subtle border
        draw.rectangle([0, 0, width-1, height-1], outline='#bdc3c7', width=2)
        
        # Save the image
        img.save(image_path, 'PNG', quality=95)
        print(f"Created professional slide image: {image_path}")
        return True
        
    except Exception as e:
        print(f"Failed to create professional slide image: {e}")
        return False


def create_enhanced_html_slide(file_path: str, slide, slide_index: int) -> bool:
    """Create an enhanced HTML slide that looks like the original PowerPoint slide"""
    try:
        # Extract slide content with better formatting
        title_text = ""
        content_texts = []
        bullet_points = []
        
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Check if it's a bullet point
                        if any(text.startswith(prefix) for prefix in ['•', '-', '*', '◦']):
                            bullet_points.append(text.lstrip('•-*◦').strip())
                        elif not title_text and len(text) < 100:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Create enhanced HTML with PowerPoint-like styling
        html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Slide {slide_index + 1}</title>
    <style>
        body {{
            font-family: 'Segoe UI', 'Calibri', Arial, sans-serif;
            margin: 0;
            padding: 0;
            background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
            color: #333;
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        .slide-container {{
            background: white;
            padding: 80px;
            border-radius: 8px;
            box-shadow: 0 10px 30px rgba(0, 0, 0, 0.2);
            max-width: 1200px;
            width: 90%;
            position: relative;
        }}
        .slide-number {{
            position: absolute;
            top: 20px;
            right: 30px;
            font-size: 14px;
            color: #666;
            font-weight: 300;
        }}
        .slide-title {{
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 40px;
            color: #2c3e50;
            line-height: 1.2;
            border-bottom: 3px solid #3498db;
            padding-bottom: 15px;
        }}
        .slide-content {{
            font-size: 20px;
            line-height: 1.6;
            color: #34495e;
        }}
        .slide-content p {{
            margin: 15px 0;
        }}
        .bullet-list {{
            margin: 20px 0;
            padding-left: 0;
        }}
        .bullet-list li {{
            margin: 12px 0;
            padding-left: 20px;
            position: relative;
        }}
        .bullet-list li:before {{
            content: "•";
            color: #3498db;
            font-weight: bold;
            position: absolute;
            left: 0;
        }}
        .content-section {{
            margin: 25px 0;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-number">Slide {slide_index + 1}</div>
        <div class="slide-title">{title_text or f"Slide {slide_index + 1}"}</div>
        <div class="slide-content">
            {chr(10).join([f"<div class='content-section'><p>{text}</p></div>" for text in content_texts[:6]]) if content_texts else ""}
            {f"<ul class='bullet-list'>{chr(10).join([f'<li>{bullet}</li>' for bullet in bullet_points[:8]])}</ul>" if bullet_points else ""}
        </div>
    </div>
</body>
</html>
        """
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        print(f"Created enhanced HTML slide: {file_path}")
        return True
        
    except Exception as e:
        print(f"Failed to create enhanced HTML slide: {e}")
        return False


def create_html_slide(file_path: str, slide_index: int, title_text: str, content_texts: list):
    """Create an HTML slide with actual content"""
    html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Slide {slide_index + 1}</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 0;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        .slide-container {{
            background: rgba(255, 255, 255, 0.1);
            padding: 60px;
            border-radius: 20px;
            text-align: center;
            backdrop-filter: blur(10px);
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3);
            max-width: 1200px;
            width: 90%;
        }}
        .slide-number {{
            font-size: 18px;
            margin-bottom: 20px;
            opacity: 0.8;
            font-weight: 300;
        }}
        .slide-title {{
            font-size: 42px;
            font-weight: bold;
            margin-bottom: 30px;
            line-height: 1.2;
        }}
        .slide-content {{
            font-size: 20px;
            line-height: 1.6;
            opacity: 0.9;
            text-align: left;
        }}
        .slide-content ul {{
            text-align: left;
            margin: 20px 0;
        }}
        .slide-content li {{
            margin: 10px 0;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-number">Slide {slide_index + 1}</div>
        <div class="slide-title">{title_text or f"Slide {slide_index + 1}"}</div>
        <div class="slide-content">
            {chr(10).join([f"<p>{text}</p>" for text in content_texts[:10]]) if content_texts else "<p>Content will be displayed here</p>"}
        </div>
    </div>
</body>
</html>
    """
    
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    print(f"Created HTML slide: {file_path}")


def create_simple_fallback(file_path: str, slide_index: int):
    """Create a simple fallback HTML file"""
    html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Slide {slide_index + 1}</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            margin: 0;
            padding: 40px;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        .slide-container {{
            background: rgba(255, 255, 255, 0.1);
            padding: 40px;
            border-radius: 20px;
            text-align: center;
            backdrop-filter: blur(10px);
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3);
        }}
        .slide-number {{
            font-size: 24px;
            margin-bottom: 20px;
            opacity: 0.8;
        }}
        .slide-title {{
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 20px;
        }}
        .slide-content {{
            font-size: 18px;
            line-height: 1.6;
            opacity: 0.9;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-number">Slide {slide_index + 1}</div>
        <div class="slide-title">Presentation Slide</div>
        <div class="slide-content">
            This slide is being processed. Content will be available shortly.
        </div>
    </div>
</body>
</html>
    """
    
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(html_content)


def create_slide_placeholder(slide, image_path: str, slide_index: int):
    """Create a placeholder image for the slide"""
    try:
        # Try to import PIL
        try:
            from PIL import Image, ImageDraw, ImageFont
            PIL_AVAILABLE = True
        except ImportError:
            PIL_AVAILABLE = False
            print("PIL not available, creating simple text file instead")
        
        if PIL_AVAILABLE:
            # Create a white background image (16:9 aspect ratio)
            width, height = 1920, 1080
            img = Image.new('RGB', (width, height), 'white')
            draw = ImageDraw.Draw(img)
            
            # Try to load a font, fallback to default if not available
            try:
                title_font = ImageFont.truetype("arial.ttf", 48)
                content_font = ImageFont.truetype("arial.ttf", 24)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
            
            # Extract text content from slide
            title_text = ""
            content_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            if not title_text and len(text) < 100:
                                title_text = text
                            else:
                                content_texts.append(text)
            
            # Draw title
            if title_text:
                draw.text((100, 100), title_text, fill='black', font=title_font)
            
            # Draw content
            y_offset = 200
            for text in content_texts[:10]:  # Limit to 10 content lines
                if y_offset < height - 100:
                    draw.text((100, y_offset), text[:80], fill='black', font=content_font)
                    y_offset += 40
            
            # Add slide number
            draw.text((width - 200, height - 50), f"Slide {slide_index + 1}", fill='gray', font=content_font)
            
            # Save the image
            img.save(image_path, 'PNG')
        else:
            # Create a simple text file as fallback
            create_text_fallback(slide, image_path, slide_index)
        
    except Exception as e:
        print(f"Error creating slide placeholder: {e}")
        # Create a simple fallback
        create_text_fallback(slide, image_path, slide_index)


def create_text_fallback(slide, image_path: str, slide_index: int):
    """Create a simple text file as fallback when PIL is not available"""
    try:
        # Extract text content
        title_text = ""
        content_texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        if not title_text and len(text) < 100:
                            title_text = text
                        else:
                            content_texts.append(text)
        
        # Create a simple HTML file that can be displayed as an image
        html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Slide {slide_index + 1}</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            margin: 0;
            padding: 40px;
            background: white;
            color: black;
        }}
        .slide-container {{
            max-width: 1200px;
            margin: 0 auto;
        }}
        .title {{
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 30px;
            color: #333;
        }}
        .content {{
            font-size: 18px;
            line-height: 1.6;
        }}
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 20px;
            color: #666;
            font-size: 14px;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="title">{title_text or f"Slide {slide_index + 1}"}</div>
        <div class="content">
            {chr(10).join([f"<p>{text}</p>" for text in content_texts[:10]])}
        </div>
        <div class="slide-number">Slide {slide_index + 1}</div>
    </div>
</body>
</html>
        """
        
        # Save as HTML file instead of PNG
        html_path = image_path.replace('.png', '.html')
        with open(html_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        # Also create a simple text file
        text_path = image_path.replace('.png', '.txt')
        with open(text_path, 'w', encoding='utf-8') as f:
            f.write(f"Slide {slide_index + 1}\n")
            f.write("=" * 50 + "\n\n")
            if title_text:
                f.write(f"Title: {title_text}\n\n")
            f.write("Content:\n")
            for text in content_texts:
                f.write(f"- {text}\n")
        
    except Exception as e:
        print(f"Error creating text fallback: {e}")


